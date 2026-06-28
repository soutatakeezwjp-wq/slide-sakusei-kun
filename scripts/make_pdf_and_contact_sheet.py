#!/usr/bin/env python3
"""Make a PDF, PPTX, and contact sheet from slide PNG/JPG images.

Usage:
  python make_pdf_and_contact_sheet.py input_dir --out output.pdf --pptx output.pptx --png-dir output_png
"""

from __future__ import annotations

import argparse
from pathlib import Path

from PIL import Image, ImageDraw

DEFAULT_BAR_COLOR = "#78a9ff"
DEFAULT_BAR_HEIGHT = 42


def list_images(input_dir: Path) -> list[Path]:
    exts = {".png", ".jpg", ".jpeg", ".webp"}
    return sorted(p for p in input_dir.iterdir() if p.suffix.lower() in exts and not p.name.startswith("contact_sheet"))


def fit_image(im: Image.Image, size: tuple[int, int]) -> Image.Image:
    im = im.convert("RGB")
    w, h = im.size
    scale = min(size[0] / w, size[1] / h)
    nw, nh = int(w * scale), int(h * scale)
    return im.resize((nw, nh), Image.LANCZOS)


def make_slide_canvas(
    im: Image.Image,
    size=(1920, 1080),
    bar_color=DEFAULT_BAR_COLOR,
    bar_height=DEFAULT_BAR_HEIGHT,
) -> Image.Image:
    """Fit the whole image inside the safe area and add mechanical bars."""
    im = im.convert("RGB")
    canvas = Image.new("RGB", size, "white")
    target_w, target_h = size
    safe_h = target_h - (bar_height * 2)
    resized = fit_image(im, (target_w, safe_h))
    canvas.paste(resized, ((target_w - resized.width) // 2, bar_height + (safe_h - resized.height) // 2))
    draw = ImageDraw.Draw(canvas)
    draw.rectangle((0, 0, target_w, bar_height), fill=bar_color)
    draw.rectangle((0, target_h - bar_height, target_w, target_h), fill=bar_color)
    return canvas


def make_pdf(images: list[Path], out_pdf: Path, bar_color: str, bar_height: int) -> None:
    pages = [make_slide_canvas(Image.open(p), bar_color=bar_color, bar_height=bar_height) for p in images]
    if not pages:
        raise SystemExit("No images found")
    pages[0].save(out_pdf, save_all=True, append_images=pages[1:], resolution=96.0)


def make_png_dir(images: list[Path], out_dir: Path, bar_color: str, bar_height: int) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    for i, path in enumerate(images, 1):
        slide = make_slide_canvas(Image.open(path), bar_color=bar_color, bar_height=bar_height)
        slide.save(out_dir / f"{i:02d}_{path.stem}.png")


def make_contact_sheet(images: list[Path], out_png: Path, bar_color: str, bar_height: int) -> None:
    thumbs = []
    for i, p in enumerate(images, 1):
        im = make_slide_canvas(Image.open(p), bar_color=bar_color, bar_height=bar_height)
        im.thumbnail((320, 180))
        cell = Image.new("RGB", (340, 230), (245, 247, 250))
        cell.paste(im, (10, 10))
        draw = ImageDraw.Draw(cell)
        draw.text((12, 198), f"{i:02d} {p.name}", fill=(0, 0, 0))
        thumbs.append(cell)
    cols = 3
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * 340, rows * 230), "white")
    for i, thumb in enumerate(thumbs):
        sheet.paste(thumb, ((i % cols) * 340, (i // cols) * 230))
    sheet.save(out_png)


def make_pptx(images: list[Path], out_pptx: Path, bar_color: str, bar_height: int) -> None:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches
    except ImportError as exc:
        raise SystemExit("python-pptx is required to create PPTX files") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    bar_h = int(slide_h * (bar_height / 1080))
    safe_h = slide_h - (bar_h * 2)

    hex_color = bar_color.lstrip("#")
    rgb = RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

    for path in images:
        slide = prs.slides.add_slide(blank)
        with Image.open(path) as im:
            w, h = im.size
        scale = min(slide_w / w, safe_h / h)
        pic_w = int(w * scale)
        pic_h = int(h * scale)
        left = int((slide_w - pic_w) / 2)
        top = int(bar_h + (safe_h - pic_h) / 2)
        slide.shapes.add_picture(str(path), left, top, width=pic_w, height=pic_h)

        for top_pos in (0, slide_h - bar_h):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, int(top_pos), slide_w, bar_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb
            shape.line.fill.background()

    prs.save(out_pptx)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("input_dir", type=Path)
    parser.add_argument("--out", type=Path, required=True)
    parser.add_argument("--png-dir", type=Path)
    parser.add_argument("--pptx", type=Path)
    parser.add_argument("--contact-sheet", type=Path)
    parser.add_argument("--bar-color", default=DEFAULT_BAR_COLOR)
    parser.add_argument("--bar-height", type=int, default=DEFAULT_BAR_HEIGHT)
    args = parser.parse_args()

    images = list_images(args.input_dir)
    if args.png_dir:
        make_png_dir(images, args.png_dir, args.bar_color, args.bar_height)
    make_pdf(images, args.out, args.bar_color, args.bar_height)
    if args.pptx:
        make_pptx(images, args.pptx, args.bar_color, args.bar_height)
    if args.contact_sheet:
        make_contact_sheet(images, args.contact_sheet, args.bar_color, args.bar_height)
    print(args.out)
    if args.png_dir:
        print(args.png_dir)
    if args.pptx:
        print(args.pptx)


if __name__ == "__main__":
    main()
